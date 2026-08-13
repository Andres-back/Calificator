import asyncio
from uuid import uuid4
from app.modules.imagenes import service as imagenes_service
from io import BytesIO

from app.modules.presentaciones import service
from app.modules.presentaciones.editable_pptx_service import build_editable_pptx
from app.modules.presentaciones.local_export import _render_slide, _role_theme
from app.modules.presentaciones.image_prompts import build_presentation_image_prompt
from app.modules.presentaciones.presentation_schema import normalize_to_canonical
from app.modules.presentaciones.schemas import PresentacionCreate
from app.shared.enums import ImageProvider


def _payload(cantidad_slides: int = 8) -> PresentacionCreate:
    return PresentacionCreate(
        titulo="Primera Ley de Newton",
        tema="Primera Ley de Newton",
        area="Fisica",
        grado="6",
        cantidad_slides=cantidad_slides,
        incluir_imagenes=True,
        densidad_imagenes="alta",
        proveedor_imagenes="premium",
    )


def test_slides_prompt_solicita_esquema_pedagogico_enriquecido() -> None:
    prompt = service.SLIDES_PROMPT

    for field in [
        "role",
        "key_message",
        "example",
        "activity",
        "question",
        "visual_concept",
        "layout_hint",
        "image_text_expected",
        "tags",
    ]:
        assert field in prompt
    for role in ["cover", "objective", "prior_knowledge", "concept", "activity", "assessment", "closing"]:
        assert role in prompt
    assert "exactitud matematica nunca depende de la imagen IA" in prompt


def test_roles_validos_se_conservan_y_roles_invalidos_tienen_fallback() -> None:
    slides = service._apply_pedagogical_slide_defaults(
        [
            {"role": "cover", "title": "Inicio", "bullets": ["Idea inicial"], "visual_concept": "Aula"},
            {"role": "rol_raro", "title": "Ejemplo guiado", "bullets": ["Idea clara"]},
        ],
        _payload(3),
    )

    assert slides[0]["role"] == "cover"
    assert slides[1]["role"] == "example"


def test_slide_intermedia_puede_ser_full_image_y_ultima_editable() -> None:
    slides = service._apply_pedagogical_slide_defaults(
        [
            {"role": "cover", "title": "Newton", "bullets": ["Inicio"], "visual_concept": "Portada"},
            {"role": "objective", "title": "Objetivo", "bullets": ["Comprender la inercia"]},
            {
                "role": "concept",
                "title": "Ley de la inercia",
                "key_message": "Los cuerpos mantienen su estado.",
                "bullets": ["Un objeto quieto permanece quieto."],
                "visual_concept": "Autobus frenando con flechas simples",
                "layout_hint": "full_image",
                "image_text_expected": ["LEY DE LA INERCIA", "Un objeto quieto permanece quieto"],
            },
            {
                "role": "activity",
                "title": "Observemos la inercia",
                "bullets": ["Retira una tarjeta rapidamente."],
                "layout_hint": "editable",
            },
        ],
        _payload(4),
    )
    legacy_full_idx = service._full_image_index(slides)

    assert service._image_kind_for_slide(slides[2], index=2, legacy_full_idx=legacy_full_idx) == "full_image"
    assert service._should_be_full_image(slides[3], index=3, legacy_full_idx=legacy_full_idx) is False
    assert slides[3]["layout_hint"] == "editable"


def test_exact_math_content_never_depends_on_generated_image() -> None:
    payload = PresentacionCreate(
        titulo="Pares de factores",
        tema="Factores y productos",
        area="Matematicas",
        grado="4",
        cantidad_slides=3,
        incluir_imagenes=True,
        densidad_imagenes="alta",
    )
    slides = service._apply_pedagogical_slide_defaults(
        [
            {
                "role": "concept",
                "title": "Dos formas de multiplicar",
                "key_message": "12 = 2 x 6",
                "bullets": ["12 = 2 x 6", "18 = 3 x 6"],
                "visual_concept": "Arreglos de circulos con cantidades exactas",
                "layout_hint": "full_image",
                "image_text_expected": ["12 = 2 x 6", "18 = 3 x 6"],
                "image_asset": "/app_data/images/incorrecta.png",
                "slide_type": "full_image",
                "layout": "full_image",
            }
        ],
        payload,
    )

    slide = slides[0]
    assert service._slide_has_exact_math_content(slide)
    assert slide["layout_hint"] == "editable"
    assert slide["image_text_expected"] == []
    assert slide["image"] == "/static/placeholder_educational.svg"
    assert "image_asset" not in slide
    assert "slide_type" not in slide
    assert "layout" not in slide
    assert service._should_be_full_image(slide, index=0, legacy_full_idx=0) is False


def test_text_only_slides_use_colorful_role_aware_layouts() -> None:
    from PIL import Image
    from pptx import Presentation

    canonical = {
        "slides": [
            {
                "tipo": "portada",
                "role": "cover",
                "layout": "cover",
                "titulo": "Factores",
                "bullets": [],
                "imagen": {},
            },
            {
                "tipo": "objetivo",
                "role": "objective",
                "layout": "text",
                "titulo": "Nuestra meta de hoy",
                "bullets": [
                    {"texto": "Reconocer factores y productos."},
                    {"texto": "Encontrar pares de factores."},
                ],
                "imagen": {},
            },
        ]
    }

    deck = Presentation(BytesIO(build_editable_pptx(canonical)))
    slide = deck.slides[1]
    fills = set()
    for shape in slide.shapes:
        try:
            rgb = shape.fill.fore_color.rgb
        except TypeError:
            continue
        if rgb is not None:
            fills.add(str(rgb))

    assert "2563EB" in fills
    assert "FFFFFF" in fills
    assert len(fills) >= 2

    rendered = Image.open(
        BytesIO(
            _render_slide(
                "Factores",
                {
                    "title": "Nuestra meta de hoy",
                    "role": "objective",
                    "bullets": [
                        "Reconocer factores y productos.",
                        "Encontrar pares de factores.",
                    ],
                },
                1,
                8,
            )
        )
    ).convert("RGB")
    assert rendered.getpixel((40, 400)) == (37, 99, 235)
    assert rendered.getpixel((1550, 400)) == (239, 246, 255)

    dark_deck_theme = ("#0b1220", "#f1f5f9", "#cbd5e1", "#94a3b8", "#a5b4fc", "#18233b")
    safe_example_theme = _role_theme("example", dark_deck_theme)
    assert safe_example_theme[0] == "#eef2ff"
    assert safe_example_theme[1] == "#1e1b4b"


def test_editable_pptx_preserves_six_dots_per_math_array_row() -> None:
    from pptx import Presentation

    dot = chr(0x25CF)
    times = chr(0x00D7)
    row = " ".join([dot] * 6)
    canonical = {
        "slides": [
            {
                "tipo": "concepto",
                "layout": "math-arrays",
                "titulo": "Arreglos exactos",
                "bullets": [
                    {"texto": f"12 = 2 {times} 6"},
                    {"texto": f"{row}\n{row}"},
                    {"texto": f"18 = 3 {times} 6"},
                    {"texto": f"{row}\n{row}\n{row}"},
                ],
                "imagen": {},
            }
        ]
    }

    deck = Presentation(BytesIO(build_editable_pptx(canonical)))
    texts = [
        paragraph.text
        for shape in deck.slides[0].shapes
        if hasattr(shape, "text_frame")
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text
    ]
    rows = [text for text in texts if dot in text]

    assert len(rows) == 5
    assert all(row_text.count(dot) == 6 for row_text in rows)


def test_activity_y_assessment_son_editables_por_defecto() -> None:
    slides = service._apply_pedagogical_slide_defaults(
        [
            {"role": "activity", "title": "Actividad", "bullets": ["Resuelve en equipo."]},
            {"role": "assessment", "title": "Evaluacion", "bullets": ["Responde la pregunta."]},
        ],
        _payload(3),
    )

    assert slides[0]["layout_hint"] == "editable"
    assert slides[1]["layout_hint"] == "editable"
    assert slides[0]["image_text_expected"] == []
    assert slides[1]["image_text_expected"] == []


def test_ocho_slides_refuerza_hitos_pedagogicos_minimos() -> None:
    slides = service._apply_pedagogical_slide_defaults(
        [
            {"role": "cover", "title": "Newton", "bullets": ["Inicio"], "visual_concept": "Portada"},
            {"role": "objective", "title": "Objetivo", "bullets": ["Comprender la inercia"]},
            {"role": "prior_knowledge", "title": "Previos", "bullets": ["Recordar movimiento"]},
            {
                "role": "concept",
                "title": "Inercia",
                "bullets": ["Los cuerpos mantienen su estado."],
                "visual_concept": "Bus frenando",
                "image_text_expected": ["INERCIA", "Los cuerpos mantienen su estado"],
            },
            {"role": "explanation", "title": "Explicacion", "bullets": ["Una fuerza cambia el estado."]},
            {"role": "example", "title": "Ejemplo", "bullets": ["Un pasajero se inclina al frenar."]},
            {"role": "activity", "title": "Actividad", "bullets": ["Observa un experimento simple."]},
            {"role": "summary", "title": "Resumen", "bullets": ["La inercia explica el movimiento."]},
        ],
        _payload(8),
    )

    assert [slide["role"] for slide in slides] == [
        "cover",
        "objective",
        "prior_knowledge",
        "concept",
        "explanation",
        "activity",
        "comprehension_check",
        "summary",
    ]
    assert slides[5]["layout_hint"] == "editable"
    assert slides[6]["layout_hint"] == "editable"
    assert slides[6]["question"]


def test_full_image_requiere_concepto_visual_y_texto_esperado() -> None:
    slide = {
        "role": "concept",
        "title": "Inercia",
        "layout_hint": "full_image",
        "visual_concept": "",
        "image_text_expected": [],
    }

    assert service._should_be_full_image(slide, index=1, legacy_full_idx=2) is False


def test_canonical_preserva_campos_pedagogicos_y_legacy_sigue_normalizando() -> None:
    canonical = normalize_to_canonical(
        {
            "title": "Newton",
            "slides": [
                {
                    "role": "concept",
                    "title": "Ley de la inercia",
                    "key_message": "Los cuerpos mantienen su estado.",
                    "bullets": ["Un objeto quieto permanece quieto."],
                    "visual_concept": "Autobus frenando",
                    "layout_hint": "full_image",
                    "image_text_expected": ["LEY DE LA INERCIA"],
                    "tags": ["fisica", "newton"],
                    "notes": "Relaciona con transporte.",
                },
                {"title": "Legacy", "bullets": ["Idea clave"], "image": "Ilustracion legacy"},
            ],
        },
        {"titulo": "Newton", "tema": "Primera Ley de Newton"},
    )

    first = canonical["slides"][0]
    second = canonical["slides"][1]
    assert first["role"] == "concept"
    assert first["key_message"] == "Los cuerpos mantienen su estado."
    assert first["visual_concept"] == "Autobus frenando"
    assert first["layout_hint"] == "full_image"
    assert first["image_text_expected"] == ["LEY DE LA INERCIA"]
    assert first["tags"] == ["fisica", "newton"]
    assert first["texto_principal"]
    assert second["title"] if "title" in second else second["titulo"]


def test_builder_full_image_usa_texto_esperado_sin_inventar_texto_visible() -> None:
    bundle = build_presentation_image_prompt(
        "full_image",
        raw_prompt="Autobus frenando con pasajeros inclinados",
        title="Ley de la inercia",
        bullets=["Este bullet no debe aparecer como texto visible."],
        topic="Primera Ley de Newton",
        area="Fisica",
        grade="6",
        provider=ImageProvider.OPENAI,
        role="concept",
        key_message="Los cuerpos mantienen su estado.",
        visual_concept="Autobus frenando con flechas de movimiento",
        image_text_expected=["LEY DE LA INERCIA", "Un objeto quieto permanece quieto"],
        tags=["fisica", "newton"],
    )

    assert bundle.prompt_original == "Autobus frenando con flechas de movimiento"
    assert '"LEY DE LA INERCIA"' in bundle.prompt_usado
    assert '"Un objeto quieto permanece quieto"' in bundle.prompt_usado
    assert "Este bullet no debe aparecer" not in bundle.prompt_usado
    assert "horizontal 16:9" in bundle.prompt_usado
    assert "sin texto pequeño" in bundle.prompt_usado
    assert bundle.image_text_expected == ["LEY DE LA INERCIA", "Un objeto quieto permanece quieto"]
    assert len(bundle.prompt_usado.split()) <= 110


def test_openai_modelo_calidad_y_hash_no_cambian() -> None:
    modelo, calidad = imagenes_service.provider_model_quality("openai")
    h1 = imagenes_service.compute_prompt_hash("abc", modelo=modelo, calidad=calidad, size=service.SLIDE_IMAGE_SIZE)
    h2 = imagenes_service.compute_prompt_hash("abc", modelo=modelo, calidad=calidad, size=service.SLIDE_IMAGE_SIZE)

    assert modelo == "gpt-image-2"
    assert calidad == "low"
    assert h1 == h2


def test_contenido_generico_del_demo_se_rechaza() -> None:
    payload = PresentacionCreate(
        titulo="Tipos de alimentacion animal",
        tema="Tipos de alimentacion animal",
        area="Ciencias Naturales",
        grado="5",
        cantidad_slides=8,
    )
    slides = service._apply_pedagogical_slide_defaults(
        service._polish_slides_for_presenton(
            service.normalize_presentation(service._fallback_slides(payload)),
            topic=payload.tema,
        ),
        payload,
    )

    issues = service._presentation_quality_issues(slides, payload)

    assert issues
    assert any("generico" in issue or "repite" in issue for issue in issues)


def test_repara_borrador_incompleto_sin_perder_la_cantidad() -> None:
    payload = PresentacionCreate(
        titulo="Fracciones equivalentes",
        tema="Fracciones equivalentes",
        area="Matematicas",
        grado="3",
        cantidad_slides=8,
    )
    draft = [
        {"title": f"Diapositiva {index + 1}", "key_message": "Idea", "bullets": ["Dato"]}
        for index in range(8)
    ]
    candidate = service._apply_pedagogical_slide_defaults(
        service.normalize_presentation(draft),
        payload,
    )

    repaired = service._repair_incomplete_presentation(candidate, payload)

    assert len(repaired) == 8
    assert service._presentation_quality_issues(repaired, payload) == []
    assert all(len(slide["bullets"]) >= 2 for slide in repaired[1:])


def test_generacion_recupera_respuesta_corta_despues_de_reintentos(monkeypatch) -> None:
    payload = PresentacionCreate(
        titulo="Fracciones equivalentes",
        tema="Fracciones equivalentes",
        area="Matematicas",
        grado="3",
        cantidad_slides=8,
    )

    class ShortContentRouter:
        calls = 0

        def __init__(self, *args, **kwargs):
            pass

        async def generate_json(self, feature, prompt):
            assert feature == "presentacion"
            self.__class__.calls += 1
            return {
                "slides": [
                    {
                        "title": f"Diapositiva {index + 1}",
                        "key_message": "Idea",
                        "bullets": ["Dato"],
                    }
                    for index in range(8)
                ]
            }

    monkeypatch.setattr(service, "LLMRouter", ShortContentRouter)

    slides = asyncio.run(service._generate_slides(payload, uuid4()))

    assert ShortContentRouter.calls == 4
    assert len(slides) == 8
    assert service._presentation_quality_issues(slides, payload) == []

def test_revision_final_corrige_un_error_conceptual() -> None:
    payload = PresentacionCreate(
        titulo="Fracciones equivalentes",
        tema="Fracciones equivalentes",
        area="Matematicas",
        grado="3",
        cantidad_slides=8,
    )
    draft = service._repair_incomplete_presentation(
        [
            {"title": f"Diapositiva {index + 1}", "bullets": ["Idea inicial"]}
            for index in range(8)
        ],
        payload,
    )
    draft[2]["bullets"][0] = "Las equivalentes tienen el mismo numerador y denominador."
    corrected = [dict(slide) for slide in draft]
    corrected[2] = {
        **corrected[2],
        "bullets": [
            "Las fracciones equivalentes representan la misma cantidad.",
            "Multiplicar numerador y denominador por el mismo numero conserva su valor.",
        ],
    }

    class Reviewer:
        async def generate_json(self, feature, prompt):
            assert feature == "presentacion"
            assert "error conceptual" in prompt
            return {"slides": corrected}

    reviewed = asyncio.run(
        service._review_slides_for_accuracy(Reviewer(), draft, payload)
    )

    visible = " ".join(reviewed[2]["bullets"]).lower()
    assert "mismo numerador y denominador" not in visible


def test_actividad_y_pregunta_quedan_visibles_en_la_exportacion() -> None:
    slides = service._merge_structured_content_into_bullets(
        [
            {
                "role": "activity",
                "bullets": ["Coloca una moneda sobre una tarjeta.", "Retira la tarjeta con rapidez."],
                "activity": "Describe por que la moneda cae dentro del vaso.",
            },
            {
                "role": "comprehension_check",
                "bullets": ["Relaciona fuerza y cambio de movimiento.", "Usa el experimento como evidencia."],
                "question": "Que propiedad explica que la moneda conserve su posicion?",
            },
        ]
    )

    assert any("Actividad:" in bullet for bullet in slides[0]["bullets"])
    assert any("Pregunta:" in bullet for bullet in slides[1]["bullets"])
