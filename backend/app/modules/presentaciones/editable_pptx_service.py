from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from app.modules.presentaciones.assets_service import resolve_asset_path


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MAX_BULLETS = 6


def build_editable_pptx(canonical: dict[str, Any]) -> bytes:
    """Build an editable PPTX from canonical presentation JSON.

    Speaker notes are intentionally not written in this phase. python-pptx does
    not expose a stable high-level API for notes pages, and skipping them keeps
    generation reliable without blocking editable slide content.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slides = (
        canonical.get("slides") if isinstance(canonical.get("slides"), list) else []
    )
    if not slides:
        slides = [_empty_slide(canonical)]

    for index, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _paint_background(slide, prs)

        layout = str(slide_data.get("layout") or "").lower()
        slide_type = str(slide_data.get("tipo") or "").lower()
        image_path = _image_path(slide_data)

        if layout == "full_image":
            if image_path:
                _add_full_image_slide(slide, slide_data, image_path)
            else:
                # Fallback editable si la imagen no está disponible.
                _add_text_slide(slide, slide_data, index=index)
        elif layout == "math-arrays":
            _add_math_arrays_slide(slide, slide_data)
        elif index == 0 or slide_type == "portada" or layout == "cover":
            _add_cover_slide(slide, slide_data, image_path)
        elif slide_type == "cierre":
            _add_closing_slide(slide, slide_data, image_path)
        elif layout in {"split-left", "split-right"} or image_path:
            _add_split_slide(
                slide, slide_data, image_path, image_left=(layout == "split-left")
            )
        else:
            _add_text_slide(slide, slide_data, index=index)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def _empty_slide(canonical: dict[str, Any]) -> dict[str, Any]:
    meta = canonical.get("meta") if isinstance(canonical.get("meta"), dict) else {}
    return {
        "tipo": "portada",
        "layout": "cover",
        "titulo": str(meta.get("titulo") or "Presentacion educativa"),
        "subtitulo": str(meta.get("tema") or ""),
        "bullets": [],
        "imagen": {},
    }


def _paint_background(slide: Any, prs: Any) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)


def _add_cover_slide(slide: Any, data: dict[str, Any], image_path: Path | None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    if image_path:
        _add_image(
            slide,
            image_path,
            Inches(7.15),
            Inches(0.0),
            Inches(6.18),
            Inches(7.5),
            crop=True,
        )
        _add_accent_bar(slide, Inches(6.95), Inches(0), Inches(0.08), Inches(7.5))
    else:
        _add_soft_panel(slide, Inches(7.15), Inches(0), Inches(6.18), Inches(7.5))

    title = _text(data.get("titulo") or "Presentacion educativa")
    subtitle = _text(data.get("subtitulo") or _first_bullet(data))
    _add_label(slide, "PRESENTACION EDUCATIVA", Inches(0.75), Inches(0.85), Inches(4.8))
    _add_textbox(
        slide,
        title,
        Inches(0.75),
        Inches(1.55),
        Inches(5.8),
        Inches(2.25),
        Pt(42),
        bold=True,
        color=RGBColor(15, 23, 42),
    )
    if subtitle:
        _add_textbox(
            slide,
            subtitle,
            Inches(0.78),
            Inches(4.15),
            Inches(5.5),
            Inches(1.0),
            Pt(20),
            color=RGBColor(71, 85, 105),
        )
    footer = _add_textbox(
        slide,
        "XCalificator",
        Inches(0.78),
        Inches(6.55),
        Inches(3),
        Inches(0.35),
        Pt(14),
        bold=True,
        color=RGBColor(13, 148, 136),
    )
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _add_text_slide(slide: Any, data: dict[str, Any], *, index: int = 0) -> None:
    """Layout editable, colorido y legible para estudiantes."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    role = str(data.get("role") or data.get("tipo") or "concept").lower()
    background, accent, soft, ink = _role_palette(role)
    _fill_slide_background(slide, background)

    bullets = _bullet_texts(data)[:4]
    title = _text(data.get("titulo"))
    label = _role_label(role)

    if role in {"activity", "actividad"}:
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.18)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = accent
        header.line.fill.background()
        _add_textbox(
            slide,
            label,
            Inches(0.85),
            Inches(0.48),
            Inches(2.5),
            Inches(0.4),
            Pt(13),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_textbox(
            slide,
            title,
            Inches(0.85),
            Inches(1.0),
            Inches(11.5),
            Inches(0.9),
            Pt(36),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_statement_tiles(
            slide,
            bullets,
            x=0.85,
            y=2.65,
            w=11.65,
            h=3.85,
            accent=accent,
            soft=soft,
            ink=ink,
            columns=2,
        )
    elif role in {"comprehension_check", "assessment", "pregunta", "evaluacion"}:
        rail = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.15), Inches(7.5)
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = accent
        rail.line.fill.background()
        question = _add_textbox(
            slide,
            "?",
            Inches(0.7),
            Inches(0.65),
            Inches(2.7),
            Inches(2.1),
            Pt(92),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        question.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _add_textbox(
            slide,
            label,
            Inches(0.75),
            Inches(3.0),
            Inches(2.9),
            Inches(0.45),
            Pt(13),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_textbox(
            slide,
            title,
            Inches(0.75),
            Inches(3.65),
            Inches(2.9),
            Inches(1.7),
            Pt(28),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_statement_tiles(
            slide,
            bullets,
            x=4.7,
            y=1.15,
            w=7.75,
            h=5.35,
            accent=accent,
            soft=soft,
            ink=ink,
            columns=1,
        )
    else:
        rail = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.78), Inches(7.5)
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = accent
        rail.line.fill.background()
        _add_textbox(
            slide,
            f"{index + 1:02d}",
            Inches(0.72),
            Inches(0.62),
            Inches(1.35),
            Inches(0.7),
            Pt(30),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_textbox(
            slide,
            label,
            Inches(0.72),
            Inches(1.55),
            Inches(2.45),
            Inches(0.5),
            Pt(13),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_textbox(
            slide,
            title,
            Inches(0.72),
            Inches(2.2),
            Inches(2.55),
            Inches(2.25),
            Pt(29),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        _add_statement_tiles(
            slide,
            bullets,
            x=4.28,
            y=1.12,
            w=8.25,
            h=5.55,
            accent=accent,
            soft=soft,
            ink=ink,
            columns=1,
        )


def _role_palette(role: str):
    from pptx.dml.color import RGBColor

    if role in {"objective", "objetivo"}:
        return (
            RGBColor(239, 246, 255),
            RGBColor(37, 99, 235),
            RGBColor(219, 234, 254),
            RGBColor(30, 58, 138),
        )
    if role in {"prior_knowledge", "saberes_previos"}:
        return (
            RGBColor(255, 247, 237),
            RGBColor(234, 88, 12),
            RGBColor(255, 237, 213),
            RGBColor(124, 45, 18),
        )
    if role in {"activity", "actividad"}:
        return (
            RGBColor(236, 253, 245),
            RGBColor(5, 150, 105),
            RGBColor(209, 250, 229),
            RGBColor(6, 78, 59),
        )
    if role in {"comprehension_check", "assessment", "pregunta", "evaluacion"}:
        return (
            RGBColor(245, 243, 255),
            RGBColor(124, 58, 237),
            RGBColor(237, 233, 254),
            RGBColor(76, 29, 149),
        )
    return (
        RGBColor(238, 242, 255),
        RGBColor(79, 70, 229),
        RGBColor(224, 231, 255),
        RGBColor(49, 46, 129),
    )


def _role_label(role: str) -> str:
    labels = {
        "objective": "NUESTRA META",
        "objetivo": "NUESTRA META",
        "prior_knowledge": "LO QUE YA SABEMOS",
        "saberes_previos": "LO QUE YA SABEMOS",
        "activity": "RETO EN CLASE",
        "actividad": "RETO EN CLASE",
        "comprehension_check": "PIENSA Y RESPONDE",
        "assessment": "DEMUESTRA LO APRENDIDO",
        "pregunta": "PIENSA Y RESPONDE",
        "evaluacion": "DEMUESTRA LO APRENDIDO",
    }
    return labels.get(role, "IDEA IMPORTANTE")


def _fill_slide_background(slide: Any, color: Any) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = color


def _add_statement_tiles(
    slide: Any,
    bullets: list[str],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    accent: Any,
    soft: Any,
    ink: Any,
    columns: int,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    visible = bullets or ["Descubre la idea principal y explicala con tus palabras."]
    columns = max(1, min(columns, 2))
    rows = (len(visible) + columns - 1) // columns
    gap = 0.28
    tile_w = (w - gap * (columns - 1)) / columns
    tile_h = min(1.55, (h - gap * max(0, rows - 1)) / max(1, rows))

    for position, text in enumerate(visible):
        column = position % columns
        row = position // columns
        tile_x = x + column * (tile_w + gap)
        tile_y = y + row * (tile_h + gap)
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(tile_x),
            Inches(tile_y),
            Inches(tile_w),
            Inches(tile_h),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        panel.line.color.rgb = soft
        panel.line.width = Pt(1.5)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(tile_x + 0.28),
            Inches(tile_y + 0.32),
            Inches(0.54),
            Inches(0.54),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        number = _add_textbox(
            slide,
            str(position + 1),
            Inches(tile_x + 0.28),
            Inches(tile_y + 0.34),
            Inches(0.54),
            Inches(0.42),
            Pt(14),
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        number.text_frame.paragraphs[0].alignment = 2
        _add_textbox(
            slide,
            _clip(text, 135),
            Inches(tile_x + 1.02),
            Inches(tile_y + 0.27),
            Inches(tile_w - 1.27),
            Inches(tile_h - 0.42),
            Pt(20),
            color=ink,
        )


def _add_math_arrays_slide(slide: Any, data: dict[str, Any]) -> None:
    """Muestra conteos como texto editable, nunca como imagen generativa."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    dot = chr(0x25CF)
    bullets = _bullet_texts(data)
    groups: list[tuple[str, list[str]]] = []
    equation = ""
    rows: list[str] = []
    for bullet in bullets:
        for line in str(bullet).splitlines():
            clean = line.strip()
            if not clean:
                continue
            if "=" in clean and not clean.startswith(dot):
                if equation:
                    groups.append((equation, rows))
                equation, rows = clean, []
            elif dot in clean:
                rows.append(clean)
    if equation:
        groups.append((equation, rows))

    _add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    _add_label(slide, "MATEMATICAS EXACTAS", Inches(0.85), Inches(0.72), Inches(2.45))
    _add_textbox(
        slide,
        _text(data.get("titulo")),
        Inches(0.85),
        Inches(1.25),
        Inches(11.6),
        Inches(0.85),
        Pt(32),
        bold=True,
        color=RGBColor(15, 23, 42),
    )

    if not groups:
        _add_bullets(
            slide,
            bullets,
            Inches(1.0),
            Inches(2.45),
            Inches(11.15),
            Inches(3.9),
            Pt(20),
        )
        return

    width = 5.55 if len(groups) > 1 else 11.15
    for index, (label, count_rows) in enumerate(groups[:2]):
        x = 0.95 + (index * 6.05)
        _add_textbox(
            slide,
            label,
            Inches(x),
            Inches(2.45),
            Inches(width),
            Inches(0.7),
            Pt(24),
            bold=True,
            color=RGBColor(13, 148, 136),
        )
        shape = slide.shapes.add_textbox(
            Inches(x),
            Inches(3.25),
            Inches(width),
            Inches(2.75),
        )
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = False
        for row_index, row in enumerate(count_rows):
            paragraph = frame.paragraphs[0] if row_index == 0 else frame.add_paragraph()
            paragraph.text = row
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(24)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(30, 41, 59)
            paragraph.space_after = Pt(10)


def _add_split_slide(
    slide: Any, data: dict[str, Any], image_path: Path | None, *, image_left: bool
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    image_w = Inches(5.2)
    if image_left:
        image_x = Inches(0)
        text_x = Inches(5.85)
    else:
        image_x = Inches(8.13)
        text_x = Inches(0.75)

    panel_x = image_x + Inches(0.35)
    panel_y = Inches(0.45)
    panel_w = image_w - Inches(0.7)
    panel_h = Inches(6.6)
    _add_soft_panel(slide, panel_x, panel_y, panel_w, panel_h)
    if image_path:
        _add_contained_image(
            slide,
            image_path,
            panel_x + Inches(0.18),
            panel_y + Inches(0.18),
            panel_w - Inches(0.36),
            panel_h - Inches(0.36),
        )

    seam_x = Inches(5.2) if image_left else Inches(8.05)
    _add_accent_bar(slide, seam_x, Inches(0), Inches(0.08), Inches(7.5))
    _add_label(slide, "DIAPOSITIVA", text_x, Inches(0.75), Inches(2.1))
    _add_textbox(
        slide,
        _text(data.get("titulo")),
        text_x,
        Inches(1.35),
        Inches(6.6),
        Inches(1.25),
        Pt(31),
        bold=True,
        color=RGBColor(15, 23, 42),
    )
    _add_bullets(
        slide,
        _bullet_texts(data),
        text_x,
        Inches(2.9),
        Inches(6.35),
        Inches(3.8),
        Pt(18),
    )


def _add_full_image_slide(slide: Any, data: dict[str, Any], image_path: Path) -> None:
    """La imagen (infografía) cubre toda la diapositiva sin deformarse:
    se escala para CUBRIR el lienzo manteniendo la proporción y se centra
    (el sobrante sangra fuera del lienzo, que PowerPoint recorta al mostrar)."""
    from pptx.util import Inches

    slide_w = SLIDE_W_IN
    slide_h = SLIDE_H_IN
    try:
        from PIL import Image as PILImage

        with PILImage.open(image_path) as img:
            img_w, img_h = img.size
        scale = max(slide_w / img_w, slide_h / img_h)
        draw_w = img_w * scale
        draw_h = img_h * scale
        x = (slide_w - draw_w) / 2
        y = (slide_h - draw_h) / 2
        slide.shapes.add_picture(
            str(image_path),
            Inches(x),
            Inches(y),
            width=Inches(draw_w),
            height=Inches(draw_h),
        )
    except Exception:  # noqa: BLE001
        # Sin PIL o imagen ilegible: ancho completo manteniendo proporción.
        slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(slide_w))


def _add_closing_slide(
    slide: Any, data: dict[str, Any], image_path: Path | None
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    if image_path:
        _add_image(
            slide,
            image_path,
            Inches(8.2),
            Inches(0.55),
            Inches(4.45),
            Inches(6.35),
            crop=True,
        )
    _add_label(slide, "CIERRE", Inches(0.9), Inches(0.8), Inches(1.5))
    title_box = _add_textbox(
        slide,
        _text(data.get("titulo") or "Cierre"),
        Inches(0.9),
        Inches(1.55),
        Inches(6.9),
        Inches(1.25),
        Pt(36),
        bold=True,
        color=RGBColor(15, 23, 42),
    )
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    message = _text(data.get("texto_principal") or _first_bullet(data))
    if message:
        _add_textbox(
            slide,
            message,
            Inches(0.95),
            Inches(3.0),
            Inches(6.6),
            Inches(1.0),
            Pt(22),
            color=RGBColor(71, 85, 105),
        )
    bullets = _bullet_texts(data)
    if len(bullets) > 1:
        _add_bullets(
            slide,
            bullets[1:],
            Inches(1.0),
            Inches(4.25),
            Inches(6.6),
            Inches(2.0),
            Pt(18),
        )


def _add_textbox(
    slide: Any,
    text: str,
    x: Any,
    y: Any,
    w: Any,
    h: Any,
    size: Any,
    *,
    bold: bool = False,
    color: Any | None = None,
) -> Any:
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = _clip(text, 180)
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Arial"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color or RGBColor(30, 41, 59)
    return shape


def _add_bullets(
    slide: Any, bullets: list[str], x: Any, y: Any, w: Any, h: Any, size: Any
) -> None:
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    visible = bullets[:MAX_BULLETS]
    if len(bullets) > MAX_BULLETS:
        visible.append("Contenido adicional resumido para mantener legibilidad.")
    if not visible:
        visible = ["Idea clave para trabajar en clase."]

    for index, bullet in enumerate(visible):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = _clip(bullet, 140)
        p.level = 0
        p.font.name = "Arial"
        p.font.size = size
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = 0


def _add_label(slide: Any, text: str, x: Any, y: Any, w: Any) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(13, 148, 136)
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def _add_accent_bar(slide: Any, x: Any, y: Any, w: Any, h: Any) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(13, 148, 136)
    shape.line.fill.background()


def _add_contained_image(
    slide: Any, path: Path, x: Any, y: Any, w: Any, h: Any
) -> None:
    if not path.is_file():
        return
    try:
        from PIL import Image as PILImage

        with PILImage.open(path) as image:
            image_w, image_h = image.size
        scale = min(float(w) / image_w, float(h) / image_h)
        draw_w = int(image_w * scale)
        draw_h = int(image_h * scale)
        draw_x = int(x + (w - draw_w) / 2)
        draw_y = int(y + (h - draw_h) / 2)
        slide.shapes.add_picture(str(path), draw_x, draw_y, width=draw_w, height=draw_h)
    except Exception:  # noqa: BLE001
        slide.shapes.add_picture(str(path), x, y, width=w)


def _add_soft_panel(slide: Any, x: Any, y: Any, w: Any, h: Any) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(226, 232, 240)
    shape.line.fill.background()


def _add_image(
    slide: Any, path: Path, x: Any, y: Any, w: Any, h: Any, *, crop: bool = False
) -> None:
    if not path.is_file():
        return
    if crop:
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    else:
        slide.shapes.add_picture(str(path), x, y, width=w)


def _image_path(slide_data: dict[str, Any]) -> Path | None:
    image = (
        slide_data.get("imagen") if isinstance(slide_data.get("imagen"), dict) else {}
    )
    raw_url = str(image.get("url") or slide_data.get("image_asset") or "").strip()
    return resolve_asset_path(raw_url)


def _bullet_texts(data: dict[str, Any]) -> list[str]:
    raw = data.get("bullets")
    if not isinstance(raw, list):
        return []
    bullets: list[str] = []
    for item in raw:
        text = str(item.get("texto") if isinstance(item, dict) else item).strip()
        if text and text != "None":
            bullets.append(text)
    return bullets


def _first_bullet(data: dict[str, Any]) -> str:
    bullets = _bullet_texts(data)
    return bullets[0] if bullets else ""


def _text(value: Any) -> str:
    return " ".join(str(value or "").split())


def _clip(value: str, max_chars: int) -> str:
    text = _text(value)
    if len(text) <= max_chars:
        return text
    return text[: max(0, max_chars - 3)].rstrip(" .,;:") + "..."
